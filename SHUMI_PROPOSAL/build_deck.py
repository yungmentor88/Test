#!/usr/bin/env python3
"""
SHUMI — Website Design Proposal deck builder.
Constructs a real, editable PPTX: live text frames, real shapes, no slide screenshots.
TedcanLabs.
"""
import os, copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT  = os.path.dirname(os.path.abspath(__file__))
NO_VIDEO = os.environ.get("NO_VIDEO") == "1"
NO_TRANS = os.environ.get("NO_TRANS") == "1"
OUTNAME  = os.environ.get("OUTNAME", "SHUMI_Website_Design_Proposal.pptx")
VID   = os.path.join(ROOT, "Higgsfield", "Videos")
IMGS  = "/home/user/Test/concepts/images"
IMGS2 = "/home/user/Test/public/images"

# ── Palette ────────────────────────────────────────────────────────────────
BONE  = RGBColor(0xF2, 0xED, 0xE6)
INK   = RGBColor(0x14, 0x10, 0x0E)
WINE  = RGBColor(0x4A, 0x12, 0x28)
ROSE  = RGBColor(0xA8, 0x26, 0x5C)
PINK  = RGBColor(0xE8, 0x5D, 0x9E)
STONE = RGBColor(0x6B, 0x62, 0x59)   # secondary text on BONE — 5.13:1
STONE_D = RGBColor(0xB3, 0xA8, 0x9E) # secondary text on INK/WINE — 8.12:1 / 6.40:1
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Office-safe faces. The production site uses Bodoni Moda + Jost; these are the
# closest pairing that ships with Office on both Windows and macOS, so the deck
# renders identically on the client's machine instead of silently substituting.
DISPLAY = "Georgia"
TEXT    = "Corbel"

W, H = Inches(13.333), Inches(7.5)

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, x, y, w, h, color, line=None, lw=1):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    s.shadow.inherit = False
    return s

def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=None, wrap=True):
    """runs: list of (string, size_pt, color, font, bold, space_after_pt, tracking_pt)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, r in enumerate(runs):
        body, size, color, font, bold, after, track = (list(r) + [None]*7)[:7]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing: p.line_spacing = line_spacing
        if after: p.space_after = Pt(after)
        run = p.add_run(); run.text = body
        f = run.font
        f.size = Pt(size); f.color.rgb = color; f.name = font or TEXT
        f.bold = bool(bold)
        if track:  # letter-spacing, in points, via raw XML
            f._rPr.set('spc', str(int(track * 100)))
    return tb

def eyebrow(slide, x, y, label, color=ROSE, w=Inches(6)):
    return text(slide, x, y, w, Inches(0.3),
                [(label.upper(), 11, color, TEXT, True, 0, 2.2)])

def rule(slide, x, y, w, color=ROSE, h=Pt(2)):
    return rect(slide, x, y, w, h, color)

def pagenum(slide, n, color=STONE):
    text(slide, W - Inches(1.0), H - Inches(0.55), Inches(0.5), Inches(0.3),
         [(f"{n:02d}", 10, color, TEXT, False, 0, 1.2)], align=PP_ALIGN.RIGHT)

def notes(slide, purpose, message, transition):
    slide.notes_slide.notes_text_frame.text = (
        f"PURPOSE — {purpose}\n\nKEY MESSAGE — {message}\n\nTRANSITION — {transition}")

def full_image(slide, path, dim=None):
    """Cover-fit an image across the whole slide, optionally with a dark scrim."""
    if not os.path.exists(path): return None
    from PIL import Image
    iw, ih = Image.open(path).size
    sa, ia = W / H, iw / ih
    if ia > sa:
        h = H; w = int(H * ia); x = int((W - w) / 2); y = 0
    else:
        w = W; h = int(W / ia); x = 0; y = int((H - h) / 2)
    pic = slide.shapes.add_picture(path, x, y, w, h)
    if dim:
        sc = rect(slide, 0, 0, W, H, INK)
        sc.fill.fore_color.rgb = INK
        sc.fill.transparency = dim
        _set_alpha(sc, dim)
    return pic

def _set_alpha(shape, alpha):
    """python-pptx has no transparency API; set it on the solidFill directly."""
    sf = shape.fill._xPr.find(qn('a:solidFill'))
    if sf is None: return
    clr = sf.find(qn('a:srgbClr'))
    if clr is None: return
    a = clr.makeelement(qn('a:alpha'), {'val': str(int((1 - alpha) * 100000))})
    clr.append(a)

def scrim(slide, x, y, w, h, color, opacity):
    s = rect(slide, x, y, w, h, color)
    _set_alpha(s, 1 - opacity)
    return s

def video(slide, path, x, y, w, h):
    """Embed a clip with its poster frame. Falls back to the poster alone."""
    poster = path.replace(".mp4", "_poster.jpg")
    if not NO_VIDEO and os.path.exists(path) and os.path.exists(poster):
        try:
            return slide.shapes.add_movie(path, x, y, w, h, poster_frame_image=poster,
                                          mime_type="video/mp4")
        except Exception:
            pass
    if os.path.exists(poster):
        return slide.shapes.add_picture(poster, x, y, w, h)
    return None

def video_full(slide, name, dim=0.55):
    """Full-bleed clip + scrim. Returns True if a clip (not just a poster) landed."""
    path = os.path.join(VID, name)
    poster = path.replace(".mp4", "_poster.jpg")
    ok = False
    if not NO_VIDEO and os.path.exists(path) and os.path.exists(poster):
        try:
            slide.shapes.add_movie(path, 0, 0, W, H, poster_frame_image=poster,
                                   mime_type="video/mp4")
            ok = True
        except Exception:
            ok = False
    if not ok and os.path.exists(poster):
        full_image(slide, poster)
    if ok or os.path.exists(poster):
        scrim(slide, 0, 0, W, H, INK, dim)
    return ok

def tag(slide, x, y, label, fill=None, fg=None, w=Inches(1.9), h=Inches(0.34), size=10):
    s = rect(slide, x, y, w, h, fill or WINE)
    text(slide, x, y + Inches(0.07), w, Inches(0.22),
         [(label.upper(), size, fg or BONE, TEXT, True, 0, 1.6)], align=PP_ALIGN.CENTER)
    return s

def placeholder(slide, x, y, w, h, label="[client to supply]", size=10):
    s = rect(slide, x, y, w, h, BONE, line=STONE, lw=0.75)
    text(slide, x, y + h/2 - Inches(0.11), w, Inches(0.22),
         [(label, size, STONE, TEXT, False, 0, 0.8)], align=PP_ALIGN.CENTER)
    return s

def morph(slide):
    """Advance this slide with a Morph transition (PowerPoint 2016+)."""
    if NO_TRANS: return
    from lxml import etree
    xml = ('<mc:AlternateContent '
           'xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006" '
           'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
           'xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" '
           'xmlns:p159="http://schemas.microsoft.com/office/powerpoint/2015/09/main">'
           '<mc:Choice Requires="p159">'
           '<p:transition spd="slow" p14:dur="1000">'
           '<p159:morph option="byObject"/>'
           '</p:transition></mc:Choice>'
           '<mc:Fallback><p:transition spd="slow"><p:fade/></p:transition></mc:Fallback>'
           '</mc:AlternateContent>')
    slide._element.append(etree.fromstring(xml))

def fade(slide, ms=700):
    """Advance this slide with a simple fade."""
    if NO_TRANS: return
    from lxml import etree
    xml = ('<mc:AlternateContent '
           'xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006" '
           'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
           'xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main">'
           '<mc:Choice Requires="p14">'
           f'<p:transition spd="slow" p14:dur="{ms}"><p:fade/></p:transition>'
           '</mc:Choice>'
           '<mc:Fallback><p:transition spd="slow"><p:fade/></p:transition></mc:Fallback>'
           '</mc:AlternateContent>')
    slide._element.append(etree.fromstring(xml))

# ══════════════════════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
M = Inches(0.9)                       # standard left margin
CW = W - M*2                          # content width

def img(*names):
    for n in names:
        for base in (IMGS, IMGS2):
            p = os.path.join(base, n)
            if os.path.exists(p): return p
    return None

# ── 01 COVER ──────────────────────────────────────────────────────────────
s = blank(prs); bg(s, INK)
video_full(s, "01_global_women.mp4", dim=0.58)
text(s, M, Inches(2.15), Inches(9.5), Inches(1.6),
     [("SHUMI", 96, BONE, DISPLAY, False, 0, -2)])
rule(s, M, Inches(3.75), Inches(1.6), PINK, Pt(3))
text(s, M, Inches(4.05), Inches(8.5), Inches(0.9),
     [("Website Design & Digital Experience", 27, BONE, DISPLAY, False, 0, 0)])
text(s, M, Inches(4.95), Inches(8.5), Inches(0.5),
     [("A global digital platform for women's empowerment", 14, PINK, TEXT, False, 0, 1.4)])
text(s, M, H - Inches(1.15), Inches(5), Inches(0.6),
     [("TEDCANLABS", 12, BONE, TEXT, True, 4, 2.4),
      ("September 2026", 11, STONE_D, TEXT, False, 0, 0.8)])
notes(s, "Open with scale and seriousness, not a logo on white.",
      "SHUMI is a global platform, and this deck behaves like one from the first frame.",
      "Hold on the footage for a beat, then: 'Before we talk about pages, let's talk about the opportunity.'")
fade(s, 900)

# ── 02 THE OPPORTUNITY ────────────────────────────────────────────────────
s = blank(prs); bg(s, BONE)
eyebrow(s, M, Inches(0.8), "01 — The opportunity")
text(s, M, Inches(1.4), Inches(7.4), Inches(2.6),
     [("A website should", 54, INK, DISPLAY, False, 2, -1.5),
      ("be more than", 54, INK, DISPLAY, False, 2, -1.5),
      ("a website.", 54, WINE, DISPLAY, False, 0, -1.5)], line_spacing=0.98)
text(s, M, Inches(4.6), Inches(5.4), Inches(1.6),
     [("SHUMI has the opportunity to build not a brochure, but a platform — "
       "one place where women, stories, programmes, events, resources, "
       "communities and partners connect to each other.", 14, STONE, TEXT, False, 0, 0)],
     line_spacing=1.5)
col = ["Women", "Stories", "Programmes", "Events", "Resources", "Communities",
       "Partners", "Opportunities"]
x0, y0 = Inches(8.6), Inches(1.5)
for i, c in enumerate(col):
    yy = y0 + Inches(0.52) * i
    rect(s, x0, yy + Inches(0.16), Inches(0.16), Pt(2), PINK)
    text(s, x0 + Inches(0.42), yy, Inches(3.2), Inches(0.4),
         [(c, 17, INK, DISPLAY, False, 0, 0)])
pagenum(s, 2)
notes(s, "Reframe the brief from 'a website' to 'a platform'.",
      "Eight content systems that currently live apart can reinforce each other.",
      "'So what should that platform feel like?'")
fade(s)

# ── 03 THE VISION ─────────────────────────────────────────────────────────
s = blank(prs); bg(s, WINE)
p = img("c1-portrait-a.jpg")
if p: s.shapes.add_picture(p, Inches(8.9), 0, Inches(4.43), H)
scrim(s, Inches(8.9), 0, Inches(4.43), H, WINE, 0.28)
eyebrow(s, M, Inches(0.8), "02 — The vision", PINK)
text(s, M, Inches(1.5), Inches(7.4), Inches(1.9),
     [("Make SHUMI's mission", 40, BONE, DISPLAY, False, 2, -1.0),
      ("experienceable.", 40, PINK, DISPLAY, False, 0, -1.0)], line_spacing=1.04)
text(s, M, Inches(3.35), Inches(6.4), Inches(1.2),
     [("Create a digital experience that makes SHUMI's mission tangible — connecting "
       "women, communities, stories and opportunities across the world.", 15, BONE, TEXT, False, 0, 0)],
     line_spacing=1.55)
for i, (k, d) in enumerate([("Global", "Not generic"),
                            ("Human", "People before jargon"),
                            ("Bold", "Confident, not loud"),
                            ("Dynamic", "Built to grow")]):
    xx = M + Inches(1.78) * i
    rule(s, xx, Inches(5.5), Inches(1.3), PINK)
    text(s, xx, Inches(5.7), Inches(1.6), Inches(1.0),
         [(k, 20, BONE, DISPLAY, False, 3, 0), (d, 10, PINK, TEXT, False, 0, 0.6)])
pagenum(s, 3, BONE)
notes(s, "State the vision in one line.",
      "The mission should be felt, not read.",
      "'Which leads to the idea the whole experience is built on.'")
fade(s)

# ── 04 THE CREATIVE IDEA ──────────────────────────────────────────────────
s = blank(prs); bg(s, INK)
eyebrow(s, M, Inches(0.85), "03 — The creative idea", PINK)
text(s, M, Inches(1.75), Inches(11.6), Inches(1.3),
     [("The world through her.", 62, BONE, DISPLAY, False, 0, -1.9)])
text(s, M, Inches(3.15), Inches(8.2), Inches(0.9),
     [("Not “women around the world” — that is a category. "
       "“The world through her” is a lens.", 15, STONE_D, TEXT, False, 0, 0)],
     line_spacing=1.5)
chain = ["HER", "HER STORY", "HER COMMUNITY", "HER OPPORTUNITY", "HER IMPACT", "THE WORLD"]
xx = M; yy = Inches(4.55)
for i, c in enumerate(chain):
    wd = Inches(1.62)
    t = rect(s, xx, yy, wd, Inches(0.6), INK, line=(PINK if i == 0 else STONE), lw=1)
    text(s, xx, yy + Inches(0.19), wd, Inches(0.3),
         [(c, 9.5, (PINK if i == 0 else BONE), TEXT, True, 0, 1.2)], align=PP_ALIGN.CENTER)
    if i < len(chain) - 1:
        text(s, xx + wd, yy + Inches(0.16), Inches(0.35), Inches(0.3),
             [("→", 13, STONE_D, TEXT, False, 0, 0)], align=PP_ALIGN.CENTER)
    xx += wd + Inches(0.35)
text(s, M, Inches(5.6), Inches(10.5), Inches(0.6),
     [("Every piece of content sits somewhere on this line. The design's job is to make the line visible.",
       13, PINK, TEXT, False, 0, 0)])
pagenum(s, 4, STONE_D)
notes(s, "Land the organising idea.",
      "One woman's story is the near end of a chain that reaches the world. That is the architecture, not a tagline.",
      "'Four principles fall out of that idea.'")
morph(s)

# ── 05 FOUR PRINCIPLES ────────────────────────────────────────────────────
s = blank(prs); bg(s, BONE)
eyebrow(s, M, Inches(0.8), "04 — Design principles")
text(s, M, Inches(1.25), Inches(8), Inches(0.8), [("Four principles.", 40, INK, DISPLAY, False, 0, -1)])
P = [("Human first", "Real people and authentic stories before organisational language.", "c2-generations.jpg"),
     ("Global by design", "A worldwide perspective carried by specificity, never by landmarks.", "c1-hero.jpg"),
     ("Impact made visible", "Turn impact into visual storytelling, not four tiles of numbers.", "c3-speaker.jpg"),
     ("Designed for movement", "An evolving platform, architected to grow without redesign.", "c3-study.jpg")]
cw = Inches(2.72); gap = Inches(0.33)
for i, (t_, d, im) in enumerate(P):
    xx = M + (cw + gap) * i
    p = img(im)
    if p: s.shapes.add_picture(p, xx, Inches(2.3), cw, Inches(1.85))
    rule(s, xx, Inches(4.35), cw, PINK)
    text(s, xx, Inches(4.55), cw, Inches(1.9),
         [(f"{i+1:02d}", 11, ROSE, TEXT, True, 6, 1.6),
          (t_, 19, INK, DISPLAY, False, 6, 0),
          (d, 11.5, STONE, TEXT, False, 0, 0)], line_spacing=1.35)
pagenum(s, 5)
notes(s, "Give the client four things to hold on to.",
      "These principles govern every later decision in the deck.",
      "'They came from studying what already works — and what does not.'")
fade(s)

# ── 06 REFERENCE LANDSCAPE ────────────────────────────────────────────────
s = blank(prs); bg(s, BONE)
eyebrow(s, M, Inches(0.8), "05 — The reference landscape")
text(s, M, Inches(1.25), Inches(9), Inches(0.9),
     [("Inspired by principles. Not copied.", 38, INK, DISPLAY, False, 0, -1)])
rect(s, M, Inches(2.35), Inches(5.5), Inches(4.1), WHITE)
rule(s, M, Inches(2.35), Inches(5.5), ROSE, Pt(3))
text(s, M + Inches(0.35), Inches(2.65), Inches(4.8), Inches(3.6),
     [("WHAT WE TAKE", 11, ROSE, TEXT, True, 12, 2),
      ("Named events as destinations", 14, INK, DISPLAY, False, 4, 0),
      ("An event with its own identity and URL converts far better than a row in a list.", 11, STONE, TEXT, False, 12, 0),
      ("Event-first architecture", 14, INK, DISPLAY, False, 4, 0),
      ("Organised around “when can I come” — the right priority for an organisation that gathers people.", 11, STONE, TEXT, False, 12, 0),
      ("Attributed testimonials", 14, INK, DISPLAY, False, 4, 0),
      ("Real names against real quotes. The credibility mechanism is right.", 11, STONE, TEXT, False, 0, 0)],
     line_spacing=1.3)
rect(s, M + Inches(5.9), Inches(2.35), Inches(5.5), Inches(4.1), INK)
rule(s, M + Inches(5.9), Inches(2.35), Inches(5.5), PINK, Pt(3))
text(s, M + Inches(6.25), Inches(2.65), Inches(4.8), Inches(3.6),
     [("WHAT WE LEAVE", 11, PINK, TEXT, True, 12, 2),
      ("The statistics band", 14, BONE, DISPLAY, False, 4, 0),
      ("8000+ attendees, 250K+ viewers. SHUMI has no such figures, and inventing them is out of the question.", 11, STONE_D, TEXT, False, 12, 0),
      ("The press-logo marquee", 14, BONE, DISPLAY, False, 4, 0),
      ("A row of logos SHUMI has not earned is a lie with a design system around it.", 11, STONE_D, TEXT, False, 12, 0),
      ("Four competing typefaces", 14, BONE, DISPLAY, False, 4, 0),
      ("And a heading order running h1 → h3 → h4 → h2, which scrambles the page for screen readers.", 11, STONE_D, TEXT, False, 0, 0)],
     line_spacing=1.3)
text(s, M, Inches(6.75), Inches(11.5), Inches(0.4),
     [("Studied from the reference's actual markup — not from impression.", 11, STONE, TEXT, False, 0, 0)])
pagenum(s, 6)
notes(s, "Show rigour, and pre-empt 'did you just copy their site?'.",
      "We analysed the real code. We take the architecture and leave the credibility scaffolding SHUMI has not earned yet.",
      "'From that, three genuinely different ways to build SHUMI.'")
fade(s)

# ── 07 THREE WAYS TO EXPERIENCE SHUMI ─────────────────────────────────────
s = blank(prs); bg(s, INK)
video_full(s, "02_human_connection.mp4", dim=0.68)
eyebrow(s, M, Inches(0.85), "06 — Three directions", PINK)
text(s, M, Inches(1.5), Inches(9), Inches(1.4),
     [("Three ways to", 50, BONE, DISPLAY, False, 2, -1.4),
      ("experience SHUMI.", 50, BONE, DISPLAY, False, 0, -1.4)], line_spacing=1.0)
D = [("01", "Editorial", "The world's editorial voice for women.", PINK),
     ("02", "Global", "A network connecting women and opportunity.", BONE),
     ("03", "Human", "A platform built on real women's stories.", BONE)]
for i, (n, t_, d, c) in enumerate(D):
    xx = M + Inches(3.85) * i
    rule(s, xx, Inches(4.5), Inches(3.2), c)
    text(s, xx, Inches(4.7), Inches(3.4), Inches(1.5),
         [(n, 12, c, TEXT, True, 6, 1.8), (t_, 30, BONE, DISPLAY, False, 6, -0.5),
          (d, 12, STONE_D, TEXT, False, 0, 0)], line_spacing=1.3)
text(s, M, Inches(6.6), Inches(11.5), Inches(0.4),
     [("Genuinely different design philosophies — not one design in three palettes.",
       12, PINK, TEXT, False, 0, 0)])
pagenum(s, 7, STONE_D)
notes(s, "Set up the three concepts.",
      "These are three different answers to what SHUMI is for, not three colourways.",
      "'Concept one.'")
fade(s)

# ── 08–10 THE THREE CONCEPTS ──────────────────────────────────────────────
CONCEPTS = [
    (8, "07", "Editorial SHUMI", WINE, BONE, PINK,
     "Credibility through craft, not statistics.",
     "SHUMI as a publication with a point of view. The only stance that lets a young "
     "organisation look serious today — with no numbers, no press and no partner logos.",
     ["Large editorial typography", "Asymmetric, grid-breaking layouts", "Full-bleed photography",
      "Oversized pull quotes", "Generous negative space"],
     [("Strength", "Degrades honestly. Works beautifully with nothing to prove."),
      ("Watch", "Coldest of the three. Lives or dies on photographic quality.")],
     ["c1-hero.jpg", "c1-portrait-b.jpg"]),
    (9, "08", "SHUMI Global", INK, BONE, PINK,
     "Geography as structure, not decoration.",
     "SHUMI as a network. Country becomes metadata on every story and programme, so "
     "'global' is a fact of the architecture rather than a claim in the copy.",
     ["Region and community storytelling", "Filtered programme discovery",
      "Impact metrics tied to stories", "Location-led navigation", "Dynamic transitions"],
     [("Strength", "Highest ceiling. The only direction that becomes a true platform."),
      ("Watch", "Requires the most content. Looks worst when empty.")],
     ["c3-hero.jpg", "c3-work.jpg"]),
    (10, "09", "Human Connection", BONE, INK, ROSE,
     "One woman, told properly.",
     "SHUMI as a deeply human platform. Documentary portraits, personal quotes and an "
     "intimate visual rhythm — story-first rather than section-first.",
     ["Large documentary portraits", "Personal quotes at scale", "Organic, warm layouts",
      "Story-first navigation", "Soft, unhurried transitions"],
     [("Strength", "Highest warmth and conversion. Best mobile experience."),
      ("Watch", "Least differentiated. Weakest at expressing 'global'.")],
     ["c2-hero.jpg", "c2-doorway.jpg"]),
]
for pg, num, title, ground, fg, accent, kicker, body, feats, tradeoffs, ims in CONCEPTS:
    s = blank(prs); bg(s, ground)
    p = img(ims[0])
    if p: s.shapes.add_picture(p, Inches(7.6), 0, Inches(5.73), H)
    scrim(s, Inches(7.6), 0, Inches(5.73), H, ground, 0.18)
    eyebrow(s, M, Inches(0.85), f"{num} — Concept {num[-1]}", accent)
    text(s, M, Inches(1.4), Inches(6.2), Inches(1.0), [(title, 42, fg, DISPLAY, False, 0, -1.2)])
    text(s, M, Inches(2.35), Inches(6.0), Inches(0.5), [(kicker, 15, accent, DISPLAY, False, 0, 0)])
    text(s, M, Inches(2.95), Inches(6.0), Inches(1.1),
         [(body, 12.5, (STONE_D if ground in (INK, WINE) else STONE), TEXT, False, 0, 0)], line_spacing=1.45)
    yy = Inches(4.15)
    for f_ in feats:
        rect(s, M, yy + Inches(0.11), Inches(0.13), Pt(1.5), accent)
        text(s, M + Inches(0.34), yy, Inches(5.6), Inches(0.3), [(f_, 12, fg, TEXT, False, 0, 0)])
        yy += Inches(0.34)
    yy += Inches(0.18)
    for lbl, d in tradeoffs:
        text(s, M, yy, Inches(6.0), Inches(0.45),
             [(lbl.upper() + " — ", 10, accent, TEXT, True, 0, 1.4), (d, 11, fg, TEXT, False, 0, 0)])
        yy += Inches(0.42)
    pagenum(s, pg, STONE_D if ground in (INK, WINE) else STONE)
    notes(s, f"Present {title} on its own terms.",
          kicker, "'And the third.'" if pg < 10 else "'So how do they compare?'")
    fade(s)

# ── 11 CONCEPT COMPARISON ─────────────────────────────────────────────────
s = blank(prs); bg(s, BONE)
eyebrow(s, M, Inches(0.75), "10 — Comparison")
text(s, M, Inches(1.15), Inches(9), Inches(0.7), [("Scored, honestly.", 36, INK, DISPLAY, False, 0, -1)])
rows = [("Premium feel", 10, 6, 8), ("Global identity", 7, 5, 10), ("Warmth", 4, 10, 5),
        ("Accessibility", 8, 10, 6), ("Storytelling", 9, 8, 7), ("Event conversion", 6, 9, 8),
        ("Mobile UX", 7, 10, 8), ("Scalability", 7, 6, 10), ("Differentiation", 9, 5, 8),
        ("Emotional impact", 7, 9, 7), ("Long-term potential", 8, 6, 10)]
tx, ty = M, Inches(2.05); rh = Inches(0.335); c0 = Inches(3.6); cw2 = Inches(1.75)
for j, h_ in enumerate(["Editorial", "Global", "Human"]):
    text(s, tx + c0 + cw2*j, ty - Inches(0.34), cw2, Inches(0.3),
         [(h_, 11, ROSE, TEXT, True, 0, 1.2)], align=PP_ALIGN.CENTER)
for i, (label, a, b, c) in enumerate(rows):
    yy = ty + rh*i
    if i % 2 == 0: rect(s, tx, yy, Inches(8.85), rh, WHITE)
    text(s, tx + Inches(0.12), yy + Inches(0.07), Inches(3.4), Inches(0.25),
         [(label, 11.5, INK, TEXT, False, 0, 0)])
    best = max(a, b, c)
    for j, v in enumerate((a, b, c)):
        win = (v == best)
        text(s, tx + c0 + cw2*j, yy + Inches(0.06), cw2, Inches(0.26),
             [(str(v), 12, ROSE if win else STONE, TEXT, win, 0, 0)], align=PP_ALIGN.CENTER)
yy = ty + rh*len(rows)
rect(s, tx, yy, Inches(8.85), Pt(2), INK)
text(s, tx + Inches(0.12), yy + Inches(0.12), Inches(3.4), Inches(0.3),
     [("Unweighted total", 12, INK, TEXT, True, 0, 0)])
for j, v in enumerate((82, 84, 87)):
    text(s, tx + c0 + cw2*j, yy + Inches(0.11), cw2, Inches(0.3),
         [(str(v), 14, INK, DISPLAY, True, 0, 0)], align=PP_ALIGN.CENTER)
rect(s, Inches(9.55), Inches(1.75), Inches(2.9), Inches(4.6), INK)
rule(s, Inches(9.55), Inches(1.75), Inches(2.9), PINK, Pt(3))
text(s, Inches(9.8), Inches(2.05), Inches(2.4), Inches(4.0),
     [("THE REAL FINDING", 10, PINK, TEXT, True, 10, 1.8),
      ("Three points across eleven categories is noise.", 13, BONE, DISPLAY, False, 10, 0),
      ("Weighted by what SHUMI needs next, each concept wins decisively — and the other two lose.",
       11, STONE_D, TEXT, False, 14, 0),
      ("Fill the event  →  Human", 11, BONE, TEXT, False, 5, 0),
      ("Win credibility  →  Editorial", 11, BONE, TEXT, False, 5, 0),
      ("Build the platform  →  Global", 11, BONE, TEXT, False, 12, 0),
      ("So the choice is a strategy decision, not a taste decision.", 11, PINK, TEXT, False, 0, 0)],
     line_spacing=1.32)
pagenum(s, 11)
notes(s, "Show the comparison is real, not rigged for a favourite.",
      "The totals are within three points. What matters is which priority SHUMI weights.",
      "'Which is why we recommend combining them.'")
fade(s)

# ── 12 RECOMMENDATION ─────────────────────────────────────────────────────
s = blank(prs); bg(s, INK)
eyebrow(s, M, Inches(0.85), "11 — Recommendation", PINK)
text(s, M, Inches(1.45), Inches(11), Inches(1.3),
     [("Editorial × Global Impact.", 56, BONE, DISPLAY, False, 0, -1.8)])
text(s, M, Inches(2.85), Inches(7.6), Inches(0.9),
     [("Editorial leads because it is the only stance that lets SHUMI look serious today. "
       "Global supplies the structure. Human supplies the heart.", 14, STONE_D, TEXT, False, 0, 0)],
     line_spacing=1.5)
mix = [("70%", "Premium editorial", "Credibility through craft", PINK),
       ("20%", "Global impact", "Geography as architecture", BONE),
       ("10%", "Human storytelling", "The emotional centre", BONE)]
for i, (pc, t_, d, c) in enumerate(mix):
    xx = M + Inches(3.85) * i
    rect(s, xx, Inches(4.15), Inches(3.35), Inches(0.09), c)
    text(s, xx, Inches(4.4), Inches(3.4), Inches(1.5),
         [(pc, 38, c, DISPLAY, False, 2, -1), (t_, 15, BONE, DISPLAY, False, 4, 0),
          (d, 11, STONE_D, TEXT, False, 0, 0)], line_spacing=1.25)
text(s, M, Inches(6.35), Inches(11.5), Inches(0.7),
     [("What we deliberately leave behind: the content debt. Counters, maps and filters are "
       "specified and switched off until SHUMI has something true to put in them.",
       11.5, PINK, TEXT, False, 0, 0)], line_spacing=1.4)
pagenum(s, 12, STONE_D)
notes(s, "Land the recommendation and its reasoning.",
      "A balance, not a formula. And we hold back the parts that would look empty.",
      "'Here is how that becomes a structure.'")
morph(s)

# ── 13 INFORMATION ARCHITECTURE ───────────────────────────────────────────
s = blank(prs); bg(s, BONE)
eyebrow(s, M, Inches(0.75), "12 — Information architecture")
text(s, M, Inches(1.15), Inches(9), Inches(0.7), [("Built to grow.", 36, INK, DISPLAY, False, 0, -1)])
IA = [("About", ["Our Story", "Mission & Vision", "Values", "Team"]),
      ("Programs", ["All Programs", "Program Detail", "Filters"]),
      ("Stories", ["Women", "Communities", "Impact"]),
      ("Global", ["Regions", "Communities", "Global Impact"]),
      ("Events", ["Upcoming", "Event Detail", "Past Events"]),
      ("Resources", ["Articles", "Reports", "Guides", "Media"]),
      ("Get Involved", ["Donate", "Partner", "Volunteer", "Participate"]),
      ("Contact", ["General", "Service inquiries"])]
hx, hy = M, Inches(2.05)
rect(s, hx, hy, Inches(1.55), Inches(0.5), INK)
text(s, hx, hy + Inches(0.15), Inches(1.55), Inches(0.3),
     [("HOME", 11, BONE, TEXT, True, 0, 1.6)], align=PP_ALIGN.CENTER)
rect(s, hx + Inches(0.74), hy + Inches(0.5), Pt(1.5), Inches(0.42), STONE)
cx = M; cy = Inches(3.0); cwid = Inches(1.42); gp = Inches(0.11)
for i, (sec, kids) in enumerate(IA):
    xx = cx + (cwid + gp) * i
    rect(s, xx, cy, cwid, Inches(0.44), WINE)
    text(s, xx, cy + Inches(0.12), cwid, Inches(0.26),
         [(sec, 9.5, BONE, TEXT, True, 0, 0.9)], align=PP_ALIGN.CENTER)
    for j, k in enumerate(kids):
        yy = cy + Inches(0.6) + Inches(0.36) * j
        rect(s, xx, yy, cwid, Inches(0.3), WHITE, line=STONE, lw=0.5)
        text(s, xx + Inches(0.06), yy + Inches(0.06), cwid - Inches(0.12), Inches(0.2),
             [(k, 8, INK, TEXT, False, 0, 0)], align=PP_ALIGN.CENTER)
rect(s, M, Inches(5.55), Inches(11.53), Inches(0.62), INK)
text(s, M + Inches(0.25), Inches(5.72), Inches(11), Inches(0.3),
     [("GROWTH SLOTS — ", 9.5, PINK, TEXT, True, 0, 1.6),
      ("Membership · Chapters · Regional hubs · Speakers · Vendors · Directory · Mentorship · Grants",
       10.5, BONE, TEXT, False, 0, 0)])
text(s, M, Inches(6.4), Inches(11.5), Inches(0.5),
     [("Each maps onto a section that already exists. Adding one never forces a navigation redesign.",
       11, STONE, TEXT, False, 0, 0)])
pagenum(s, 13)
notes(s, "Show the structure is designed for 2028, not just launch.",
      "Nine sections, and eight future features that already have a home.",
      "'Now the homepage as a journey through that structure.'")
morph(s)

# ── 14 HOMEPAGE JOURNEY ───────────────────────────────────────────────────
s = blank(prs); bg(s, INK)
eyebrow(s, M, Inches(0.8), "13 — Homepage journey", PINK)
text(s, M, Inches(1.2), Inches(9), Inches(0.8),
     [("A journey, not a stack of sections.", 34, BONE, DISPLAY, False, 0, -0.9)])
J = [("01", "Hero"), ("02", "Introduction"), ("03", "Global impact"), ("04", "Stories"),
     ("05", "Programs"), ("06", "Global connection"), ("07", "Events"), ("08", "Resources"),
     ("09", "Partners"), ("10", "Get involved"), ("11", "Newsletter"), ("12", "Footer")]
bx, by = M, Inches(2.5); bw = Inches(0.86); bgap = Inches(0.11)
for i, (n, lab) in enumerate(J):
    xx = bx + (bw + bgap) * i
    hgt = Inches(1.5) if i in (0, 2, 3, 6) else Inches(0.95)
    c = PINK if i in (0, 2, 3, 6) else WINE
    rect(s, xx, by + (Inches(1.5) - hgt), bw, hgt, c)
    text(s, xx, by + Inches(1.62), bw, Inches(0.25),
         [(n, 8.5, STONE_D, TEXT, True, 0, 1)], align=PP_ALIGN.CENTER)
    tb = text(s, xx - Inches(0.18), by + Inches(1.92), bw + Inches(0.36), Inches(0.7),
              [(lab, 8.5, BONE, TEXT, False, 0, 0)], align=PP_ALIGN.CENTER)
text(s, M, Inches(5.35), Inches(11.5), Inches(0.4),
     [("Highlighted: the four beats that carry the argument.", 10.5, PINK, TEXT, False, 0, 0)])
rect(s, M, Inches(5.95), Inches(11.53), Inches(0.95), WINE)
text(s, M + Inches(0.3), Inches(6.15), Inches(11), Inches(0.6),
     [("THE TRANSITION RULE", 9.5, PINK, TEXT, True, 5, 1.6),
      ("Consecutive sections never share a ground colour or a layout shape. That rhythm is what stops twelve sections reading as twelve boxes.",
       11.5, BONE, TEXT, False, 0, 0)])
pagenum(s, 14, STONE_D)
notes(s, "Show the homepage is composed, not assembled.",
      "Each section hands off to the next. The rhythm is designed.",
      "'It opens here.'")
morph(s)

# ── 15 HERO EXPERIENCE ────────────────────────────────────────────────────
s = blank(prs); bg(s, INK)
video_full(s, "03_homepage_hero.mp4", dim=0.5)
text(s, M, Inches(1.9), Inches(8.6), Inches(2.4),
     [("When women move forward,", 40, BONE, DISPLAY, False, 4, -1.1),
      ("the world moves with them.", 40, BONE, DISPLAY, False, 0, -1.1)], line_spacing=1.02)
text(s, M, Inches(4.15), Inches(6.2), Inches(0.6),
     [("SHUMI connects women across countries and generations.", 14, BONE, TEXT, False, 0, 0)])
b1 = rect(s, M, Inches(5.0), Inches(2.0), Inches(0.56), PINK)
text(s, M, Inches(5.17), Inches(2.0), Inches(0.3),
     [("EXPLORE SHUMI", 10.5, INK, TEXT, True, 0, 1.4)], align=PP_ALIGN.CENTER)
b2 = rect(s, M + Inches(2.25), Inches(5.0), Inches(2.0), Inches(0.56), INK, line=BONE, lw=1)
text(s, M + Inches(2.25), Inches(5.17), Inches(2.0), Inches(0.3),
     [("MEET THE WOMEN", 10.5, BONE, TEXT, True, 0, 1.4)], align=PP_ALIGN.CENTER)
eyebrow(s, M, Inches(0.85), "14 — Hero experience", PINK)
rect(s, Inches(9.05), Inches(5.72), Inches(3.38), Inches(1.02), INK)
rule(s, Inches(9.05), Inches(5.72), Inches(3.38), PINK, Pt(2))
text(s, Inches(9.25), Inches(5.92), Inches(3.0), Inches(0.8),
     [("Composition briefed with the headline in mind — the type never needs a scrim to be legible.",
       10.5, BONE, TEXT, False, 0, 0)], line_spacing=1.3)
pagenum(s, 15, STONE)
notes(s, "Show the hero working, not described.",
      "Headline is a creative direction, not approved copy. The negative space is a hard requirement.",
      "'Then SHUMI introduces itself.'")
fade(s)

# ── 16 SHUMI INTRODUCTION ─────────────────────────────────────────────────
s = blank(prs); bg(s, BONE)
p = img("c1-group.jpg")
if p: s.shapes.add_picture(p, 0, 0, Inches(5.4), H)
eyebrow(s, Inches(6.2), Inches(1.0), "15 — Introduction")
text(s, Inches(6.2), Inches(1.5), Inches(6.2), Inches(2.8),
     [("A manifesto.", 40, INK, DISPLAY, False, 8, -1),
      ("Not a paragraph.", 40, WINE, DISPLAY, False, 0, -1)], line_spacing=1.0)
text(s, Inches(6.2), Inches(3.65), Inches(5.9), Inches(1.4),
     [("The introduction is set as a large typographic statement, broken across lines and "
       "revealed as the visitor arrives. Short, declarative, in SHUMI's own voice.",
       13, STONE, TEXT, False, 0, 0)], line_spacing=1.5)
rect(s, Inches(6.2), Inches(5.05), Inches(5.9), Inches(1.5), WHITE)
rule(s, Inches(6.2), Inches(5.05), Inches(5.9), ROSE, Pt(3))
text(s, Inches(6.45), Inches(5.32), Inches(5.4), Inches(1.1),
     [("“[SHUMI'S OWN WORDS — 25–40 words, to be written with SHUMI]”", 15, INK, DISPLAY, False, 8, 0),
      ("Never fabricated. This slot stays visibly empty until SHUMI writes it.", 10.5, STONE, TEXT, False, 0, 0)],
     line_spacing=1.3)
pagenum(s, 16)
notes(s, "Show restraint as a design decision.",
      "We are not writing SHUMI's mission for them. The component is built; the words are theirs.",
      "'Then impact — the part everyone gets wrong.'")
fade(s)

# ── 17 IMPACT EXPERIENCE ──────────────────────────────────────────────────
s = blank(prs); bg(s, INK)
video_full(s, "04_global_impact.mp4", dim=0.72)
eyebrow(s, M, Inches(0.8), "16 — Impact experience", PINK)
text(s, M, Inches(1.2), Inches(9), Inches(0.9),
     [("A number that hands you a person.", 34, BONE, DISPLAY, False, 0, -0.9)])
stats = [("XX,XXX", "Women reached"), ("XX", "Communities"), ("XX", "Countries"), ("XX", "Programmes")]
for i, (v, l) in enumerate(stats):
    xx = M + Inches(2.9) * i
    text(s, xx, Inches(2.35), Inches(2.7), Inches(1.0),
         [(v, 42, PINK, DISPLAY, False, 2, -1.2), (l, 11, BONE, TEXT, False, 0, 0.8)])
rule(s, M, Inches(3.75), Inches(11.53), STONE_D, Pt(1))
text(s, M, Inches(4.0), Inches(5.6), Inches(1.9),
     [("Then the chain continues.", 17, BONE, DISPLAY, False, 8, 0),
      ("A statistic that hands straight to one woman is worth more than four statistics in tiles. "
       "It also degrades gracefully — with the number still a placeholder, her story alone carries the section.",
       12, STONE_D, TEXT, False, 0, 0)], line_spacing=1.45)
chain2 = ["XX,XXX WOMEN", "ONE WOMAN", "HER COMMUNITY", "THE WORLD"]
xx = Inches(7.0)
for i, c in enumerate(chain2):
    rect(s, xx, Inches(4.35), Inches(1.15), Inches(0.55), INK, line=(PINK if i == 1 else STONE), lw=1)
    text(s, xx + Inches(0.04), Inches(4.5), Inches(1.07), Inches(0.5),
         [(c, 7.5, (PINK if i == 1 else BONE), TEXT, True, 0, 0.6)], align=PP_ALIGN.CENTER)
    if i < 3:
        text(s, xx + Inches(1.15), Inches(4.47), Inches(0.28), Inches(0.3),
             [("→", 11, STONE_D, TEXT, False, 0, 0)], align=PP_ALIGN.CENTER)
    xx += Inches(1.43)
rect(s, M, Inches(6.2), Inches(11.53), Inches(0.62), WINE)
text(s, M + Inches(0.28), Inches(6.38), Inches(11), Inches(0.3),
     [("ALL FIGURES ARE PLACEHOLDERS. ", 9.5, PINK, TEXT, True, 0, 1.5),
      ("No SHUMI statistic has been invented. These slots stay visibly empty until real numbers exist.",
       10.5, BONE, TEXT, False, 0, 0)])
pagenum(s, 17, STONE_D)
notes(s, "Show the impact mechanism and the honesty policy at once.",
      "Impact is not four tiles. It is a chain from a number to a person.",
      "'And the same chain runs outward, geographically.'")
fade(s)

# ── 18 GLOBAL EXPERIENCE ──────────────────────────────────────────────────
s = blank(prs); bg(s, BONE)
eyebrow(s, M, Inches(0.75), "17 — Global experience")
text(s, M, Inches(1.15), Inches(9.5), Inches(0.8),
     [("Explorable. Not a spinning globe.", 34, INK, DISPLAY, False, 0, -0.9)])
steps = ["WORLD", "REGION", "COMMUNITY", "PROGRAM", "WOMAN", "STORY"]
xx = M
for i, st in enumerate(steps):
    wd = Inches(1.72)
    c = WINE if i < 3 else ROSE
    rect(s, xx, Inches(2.15), wd, Inches(0.52), c)
    text(s, xx, Inches(2.3), wd, Inches(0.3),
         [(st, 9.5, BONE, TEXT, True, 0, 1.2)], align=PP_ALIGN.CENTER)
    if i < 5:
        text(s, xx + wd, Inches(2.28), Inches(0.24), Inches(0.3),
             [("›", 14, STONE, TEXT, True, 0, 0)], align=PP_ALIGN.CENTER)
    xx += wd + Inches(0.24)
regions = ["Africa", "Americas", "Asia", "Europe", "Middle East", "Oceania"]
rect(s, M, Inches(3.15), Inches(5.4), Inches(3.1), WHITE)
text(s, M + Inches(0.28), Inches(3.38), Inches(4.9), Inches(0.3),
     [("FILTER BY REGION", 9.5, ROSE, TEXT, True, 0, 1.6)])
for i, r in enumerate(regions):
    rx = M + Inches(0.28) + (Inches(1.62) * (i % 3))
    ry = Inches(3.75) + Inches(0.5) * (i // 3)
    rect(s, rx, ry, Inches(1.5), Inches(0.38), BONE if i else WINE, line=STONE if i else None, lw=0.75)
    text(s, rx, ry + Inches(0.09), Inches(1.5), Inches(0.24),
         [(r, 9, BONE if not i else INK, TEXT, i == 0, 0, 0)], align=PP_ALIGN.CENTER)
placeholder(s, M + Inches(0.28), Inches(4.85), Inches(4.85), Inches(0.4),
            "[COUNTRIES SHUMI OPERATES IN — client to supply]")
rect(s, Inches(6.7), Inches(3.15), Inches(5.73), Inches(3.1), INK)
rule(s, Inches(6.7), Inches(3.15), Inches(5.73), PINK, Pt(3))
text(s, Inches(6.98), Inches(3.45), Inches(5.2), Inches(2.6),
     [("WHY A LIST BEFORE A MAP", 9.5, PINK, TEXT, True, 10, 1.6),
      ("A typeset region index is honest at three countries and still works at sixty. "
       "A world map with four pins advertises how few pins there are — the opposite of the intended effect.",
       12, BONE, TEXT, False, 12, 0),
      ("It is also keyboard-navigable and legible to a screen reader, which a map is not without "
       "significant extra work. The map becomes progressive enhancement, never the only way in.",
       11, STONE_D, TEXT, False, 0, 0)], line_spacing=1.4)
pagenum(s, 18)
notes(s, "Pre-empt the client asking for an animated globe.",
      "We can build the map. We recommend earning it first.",
      "'Programmes work the same way.'")
fade(s)

# ── 19 PROGRAM EXPERIENCE ─────────────────────────────────────────────────
s = blank(prs); bg(s, BONE)
eyebrow(s, M, Inches(0.75), "18 — Programme experience")
text(s, M, Inches(1.15), Inches(9), Inches(0.7),
     [("Discovery that works at six, and at sixty.", 32, INK, DISPLAY, False, 0, -0.9)])
filters = ["All", "Region", "Theme", "Type", "Status"]
fx = M
for i, f_ in enumerate(filters):
    wd = Inches(1.15)
    rect(s, fx, Inches(2.1), wd, Inches(0.4), WINE if i == 0 else BONE, line=None if i == 0 else STONE, lw=0.75)
    text(s, fx, Inches(2.2), wd, Inches(0.24),
         [(f_, 9.5, BONE if i == 0 else INK, TEXT, i == 0, 0, 0.8)], align=PP_ALIGN.CENTER)
    fx += wd + Inches(0.14)
cards = [("c3-work.jpg", "EDUCATION"), ("c2-hands.jpg", "ECONOMIC"),
         ("c3-study.jpg", "LEADERSHIP"), ("c1-portrait-b.jpg", "COMMUNITY")]
cw3 = Inches(2.72); gp3 = Inches(0.33)
for i, (im, cat) in enumerate(cards):
    xx = M + (cw3 + gp3) * i
    p = img(im)
    if p: s.shapes.add_picture(p, xx, Inches(2.85), cw3, Inches(1.75))
    text(s, xx, Inches(4.75), cw3, Inches(2.0),
         [(cat, 9, ROSE, TEXT, True, 6, 1.6),
          ("[PROGRAM NAME]", 17, INK, DISPLAY, False, 5, 0),
          ("[Region] · [Theme]", 10, STONE, TEXT, False, 6, 0),
          ("[Short description — client to supply]", 10.5, STONE, TEXT, False, 8, 0),
          ("Explore →", 10.5, ROSE, TEXT, True, 0, 0)], line_spacing=1.3)
pagenum(s, 19)
notes(s, "Show a scalable content system, not four hand-built pages.",
      "One component, filtered. It looks right with three programmes and with thirty.",
      "'The same discipline drives the most important system — stories.'")
fade(s)

# ── 20 STORIES EXPERIENCE ─────────────────────────────────────────────────
s = blank(prs); bg(s, WINE)
has = video_full(s, "05_stories.mp4", dim=0.62)
if not has and not os.path.exists(os.path.join(VID, "05_stories_poster.jpg")):
    p = img("c1-portrait-a.jpg")
    if p: s.shapes.add_picture(p, Inches(7.9), 0, Inches(5.43), H)
    scrim(s, 0, 0, W, H, WINE, 0.35)
eyebrow(s, M, Inches(0.8), "19 — Stories experience", PINK)
text(s, M, Inches(1.25), Inches(8), Inches(0.9),
     [("Her story is part of a larger story.", 34, BONE, DISPLAY, False, 0, -0.9)])
rect(s, M, Inches(2.35), Inches(6.1), Inches(3.9), INK)
rule(s, M, Inches(2.35), Inches(6.1), PINK, Pt(3))
text(s, M + Inches(0.38), Inches(2.7), Inches(5.4), Inches(3.4),
     [("[WOMAN'S NAME]", 26, BONE, DISPLAY, False, 5, -0.6),
      ("[LOCATION] · [REGION]", 10, PINK, TEXT, True, 12, 1.6),
      ("“[Her quotation — real, short, attributed, and used with her permission.]”",
       15, BONE, DISPLAY, False, 12, 0),
      ("[Short introduction — 40 to 60 words, client to supply]", 11, STONE_D, TEXT, False, 14, 0),
      ("READ HER STORY  →", 10.5, PINK, TEXT, True, 0, 1.4)], line_spacing=1.35)
text(s, Inches(7.6), Inches(5.05), Inches(4.9), Inches(1.5),
     [("Every story carries region, theme and programme as metadata. That is what makes the "
       "global architecture real rather than claimed — and what lets one story link back up "
       "to its community, its region and the world.", 11.5, BONE, TEXT, False, 0, 0)], line_spacing=1.45)
pagenum(s, 20, STONE_D)
notes(s, "This is the emotional centre of the deck. Slow down here.",
      "One woman, told properly, beats six cards. And her record is structured so it connects outward.",
      "'Stories bring people in. Events bring them together.'")
fade(s)

# ── 21 EVENTS EXPERIENCE ──────────────────────────────────────────────────
s = blank(prs); bg(s, BONE)
eyebrow(s, M, Inches(0.75), "20 — Events experience")
text(s, M, Inches(1.15), Inches(9), Inches(0.7),
     [("Events are experiences, not a calendar.", 32, INK, DISPLAY, False, 0, -0.9)])
p = img("c1-group.jpg")
if p: s.shapes.add_picture(p, M, Inches(2.15), Inches(5.4), Inches(3.05))
rect(s, Inches(6.7), Inches(2.15), Inches(5.73), Inches(3.05), INK)
text(s, Inches(7.0), Inches(2.45), Inches(5.1), Inches(2.6),
     [("[DATE]", 11, PINK, TEXT, True, 8, 1.8),
      ("[EVENT NAME]", 30, BONE, DISPLAY, False, 8, -0.8),
      ("[LOCATION]", 11, STONE, TEXT, False, 10, 1),
      ("[Short description — client to supply]", 11.5, STONE, TEXT, False, 14, 0)], line_spacing=1.35)
rect(s, Inches(7.0), Inches(4.5), Inches(2.1), Inches(0.5), PINK)
text(s, Inches(7.0), Inches(4.64), Inches(2.1), Inches(0.3),
     [("REGISTER", 10, INK, TEXT, True, 0, 1.5)], align=PP_ALIGN.CENTER)
rect(s, M, Inches(5.45), Inches(11.53), Inches(1.15), WHITE)
rule(s, M, Inches(5.45), Inches(11.53), ROSE, Pt(3))
text(s, M + Inches(0.3), Inches(5.7), Inches(11), Inches(0.9),
     [("ARCHITECTED FOR EVENTBRITE — NOT INTEGRATED TODAY", 9.5, ROSE, TEXT, True, 7, 1.6),
      ("The ticket action is a slot behind a clean interface, so connecting Eventbrite later is "
       "configuration rather than a rebuild. We are not claiming an integration that does not exist. "
       "A post-event state is designed too, so the section never looks stale the morning after.",
       11.5, STONE, TEXT, False, 0, 0)], line_spacing=1.4)
pagenum(s, 21)
notes(s, "Named events as destinations — the one idea worth taking from the reference.",
      "And an honest position on Eventbrite: architected for, not claimed.",
      "'Once someone cares, they need a way in.'")
fade(s)

# ── 22 GET INVOLVED ───────────────────────────────────────────────────────
s = blank(prs); bg(s, INK)
eyebrow(s, M, Inches(0.8), "21 — Get involved", PINK)
text(s, M, Inches(1.2), Inches(9), Inches(0.8),
     [("Learn → Connect → Act.", 34, BONE, DISPLAY, False, 0, -0.9)])
text(s, M, Inches(2.05), Inches(8.4), Inches(0.5),
     [("Not six identical buttons. A progression ordered by commitment.", 13, STONE_D, TEXT, False, 0, 0)])
paths = [("Share", "Lowest effort. One tap.", 1), ("Attend", "Come to an event.", 2),
         ("Volunteer", "Give time.", 3), ("Participate", "Join a programme.", 4),
         ("Partner", "Organisational.", 5), ("Donate", "Highest commitment.", 6)]
bx = M
for i, (t_, d, lvl) in enumerate(paths):
    wd = Inches(1.78)
    hgt = Inches(0.55) + Inches(0.19) * lvl
    yy = Inches(5.15) - hgt
    c = PINK if i >= 4 else WINE
    rect(s, bx, yy, wd, hgt, c)
    text(s, bx + Inches(0.14), yy + Inches(0.14), wd - Inches(0.28), Inches(0.5),
         [(t_, 15, BONE if i < 4 else INK, DISPLAY, False, 0, 0)])
    text(s, bx, Inches(5.3), wd, Inches(0.6),
         [(d, 9.5, STONE_D, TEXT, False, 0, 0)])
    bx += wd + Inches(0.17)
text(s, M, Inches(6.15), Inches(11.5), Inches(0.7),
     [("Height signals commitment. The page moves left to right from the easiest thing a visitor "
       "can do to the biggest — so nobody is asked to donate before they have been given a reason to.",
       11.5, PINK, TEXT, False, 0, 0)], line_spacing=1.4)
pagenum(s, 22, STONE_D)
notes(s, "Show conversion thinking, not a button row.",
      "The ask escalates. Donation is the end of a journey, not the front door.",
      "'And most of this happens on a phone.'")
fade(s)

# ── 23 MOBILE EXPERIENCE ──────────────────────────────────────────────────
s = blank(prs); bg(s, BONE)
eyebrow(s, M, Inches(0.75), "22 — Mobile experience")
text(s, M, Inches(1.15), Inches(9), Inches(0.7),
     [("Designed for the phone, not shrunk to fit.", 32, INK, DISPLAY, False, 0, -0.9)])
# three device frames
dev = [("Home", "c2-hero.jpg"), ("Story", "c1-portrait-a.jpg"), ("Event", "c1-group.jpg")]
for i, (lab, im) in enumerate(dev):
    dx = M + Inches(1.95) * i
    rect(s, dx, Inches(2.15), Inches(1.62), Inches(3.4), INK)
    p = img(im)
    if p: s.shapes.add_picture(p, dx + Inches(0.07), Inches(2.22), Inches(1.48), Inches(1.5))
    rect(s, dx + Inches(0.07), Inches(3.8), Inches(1.48), Inches(0.09), PINK)
    for k in range(4):
        rect(s, dx + Inches(0.07), Inches(3.98) + Inches(0.16)*k,
             Inches(1.48) if k % 2 == 0 else Inches(1.0), Inches(0.07), STONE)
    rect(s, dx + Inches(0.07), Inches(5.18), Inches(1.48), Inches(0.28), PINK)
    text(s, dx, Inches(5.62), Inches(1.62), Inches(0.3),
         [(lab, 10, STONE, TEXT, True, 0, 1.2)], align=PP_ALIGN.CENTER)
rules = [("Menu says the word “Menu”", "Not a bare icon. Older users tap words."),
         ("Primary action in the lower third", "Thumb reach is a layout constraint."),
         ("Art-directed crops per breakpoint", "A 21:9 desktop hero becomes 4:5 on mobile."),
         ("Nothing depends on hover", "Every menu opens on tap."),
         ("Scale relationships hold", "Display steps 112px → 44px; the ratios survive.")]
for i, (t_, d) in enumerate(rules):
    yy = Inches(2.2) + Inches(0.72) * i
    rect(s, Inches(6.7), yy + Inches(0.09), Inches(0.13), Pt(1.5), ROSE)
    text(s, Inches(6.98), yy, Inches(5.4), Inches(0.65),
         [(t_, 13, INK, DISPLAY, False, 3, 0), (d, 10.5, STONE, TEXT, False, 0, 0)], line_spacing=1.3)
pagenum(s, 23)
notes(s, "Mobile is where most visitors actually arrive.",
      "These are decisions, not responsive defaults.",
      "'Motion is the last layer of the experience.'")
fade(s)

# ── 24 MOTION SYSTEM ──────────────────────────────────────────────────────
s = blank(prs); bg(s, INK)
eyebrow(s, M, Inches(0.8), "23 — Motion system", PINK)
text(s, M, Inches(1.2), Inches(9), Inches(0.8),
     [("Motion with a job to do.", 34, BONE, DISPLAY, False, 0, -0.9)])
mot = [("Image reveal", "Mask wipe up, 600ms. Once, on first view."),
       ("Editorial statement", "Line-by-line rise, 80ms stagger."),
       ("Impact figures", "Count up 1200ms — and present as static text before any script runs."),
       ("Story rail", "Horizontal momentum, keyboard operable."),
       ("Page transition", "300ms cross-fade."),
       ("Hover", "150ms colour only. Never scale, never lift.")]
for i, (t_, d) in enumerate(mot):
    xx = M + Inches(3.9) * (i % 3)
    yy = Inches(2.3) + Inches(1.35) * (i // 3)
    rule(s, xx, yy, Inches(3.3), PINK)
    text(s, xx, yy + Inches(0.18), Inches(3.4), Inches(1.0),
         [(t_, 15, BONE, DISPLAY, False, 5, 0), (d, 10.5, STONE_D, TEXT, False, 0, 0)], line_spacing=1.3)
rect(s, M, Inches(5.35), Inches(11.53), Inches(1.35), WINE)
text(s, M + Inches(0.3), Inches(5.6), Inches(11), Inches(1.1),
     [("FORBIDDEN", 9.5, PINK, TEXT, True, 7, 1.6),
      ("Bounce · spin · fly-in · auto-advancing carousels · scroll-jacking · parallax on text · "
       "anything that repeats on every scroll past.", 12, BONE, TEXT, False, 8, 0),
      ("prefers-reduced-motion disables all of it. Nothing is discoverable only through motion, and no "
       "information exists only inside an animation.", 11, STONE_D, TEXT, False, 0, 0)], line_spacing=1.35)
pagenum(s, 24, STONE_D)
notes(s, "Show motion is governed, not sprinkled.",
      "Six moments, each with a reason. And a reduced-motion path that actually works.",
      "'Underneath all of it sits the content system.'")
fade(s)

# ── 25 CONTENT SYSTEM ─────────────────────────────────────────────────────
s = blank(prs); bg(s, BONE)
eyebrow(s, M, Inches(0.75), "24 — Content system")
text(s, M, Inches(1.15), Inches(9), Inches(0.7),
     [("SHUMI updates the site. Not a developer.", 32, INK, DISPLAY, False, 0, -0.9)])
ents = [("Story", "Name · Location · Portrait · Quote · Body"),
        ("Program", "Name · Category · Impact · Status"),
        ("Event", "Name · Date · Location · Ticket URL"),
        ("Resource", "Title · Type · File · Description"),
        ("Region", "Name · Communities"),
        ("Impact metric", "Label · Value · Source · As-of date"),
        ("Partner", "Name · Logo · Consent flag"),
        ("Team", "Name · Role · Bio · Consent flag")]
for i, (t_, f_) in enumerate(ents):
    xx = M + Inches(2.9) * (i % 4)
    yy = Inches(2.15) + Inches(1.12) * (i // 4)
    rect(s, xx, yy, Inches(2.7), Inches(0.92), WHITE)
    rule(s, xx, yy, Inches(2.7), WINE, Pt(2.5))
    text(s, xx + Inches(0.2), yy + Inches(0.18), Inches(2.4), Inches(0.65),
         [(t_, 15, INK, DISPLAY, False, 4, 0), (f_, 9, STONE, TEXT, False, 0, 0)], line_spacing=1.25)
rect(s, M, Inches(4.65), Inches(11.53), Inches(1.55), INK)
rule(s, M, Inches(4.65), Inches(11.53), PINK, Pt(3))
text(s, M + Inches(0.32), Inches(4.95), Inches(11), Inches(1.2),
     [("TWO FIELDS THAT ARE NOT DECORATION", 9.5, PINK, TEXT, True, 8, 1.6),
      ("Impact metric carries Source and As-of date. Partner and Team carry a Consent flag.",
       13, BONE, DISPLAY, False, 8, 0),
      ("These make it structurally awkward to publish a number nobody can stand behind, or to name a "
       "person or organisation that has not agreed. The content model enforces the honesty policy — "
       "rather than relying on whoever happens to be editing that day.",
       11, STONE_D, TEXT, False, 0, 0)], line_spacing=1.4)
pagenum(s, 25)
notes(s, "Show the CMS is designed, and that it protects SHUMI.",
      "Nine entities, relational. Two fields exist specifically to prevent accidental fabrication.",
      "'Briefly, what it runs on.'")
fade(s)

# ── 26 TECHNICAL FOUNDATION ───────────────────────────────────────────────
s = blank(prs); bg(s, INK)
eyebrow(s, M, Inches(0.8), "25 — Technical foundation", PINK)
text(s, M, Inches(1.2), Inches(9), Inches(0.8),
     [("Chosen for reasons, not résumé.", 34, BONE, DISPLAY, False, 0, -0.9)])
tech = [("Next.js + TypeScript", "Editorial sites live or die on image performance and SEO."),
        ("Headless CMS", "The content model is relational; a page-based CMS cannot express it."),
        ("Tailwind + design tokens", "Tokens map 1:1 to the system, so contrast rules survive handoff."),
        ("Art-directed images", "AVIF/WebP, per breakpoint. Cinematic must not mean slow."),
        ("Accessibility built in", "WCAG 2.1 AA measured, not assumed."),
        ("Eventbrite adapter", "Behind a clean interface. Swappable, and not claimed as built.")]
for i, (t_, d) in enumerate(tech):
    xx = M + Inches(3.9) * (i % 3)
    yy = Inches(2.3) + Inches(1.3) * (i // 3)
    rule(s, xx, yy, Inches(3.3), PINK)
    text(s, xx, yy + Inches(0.18), Inches(3.45), Inches(1.0),
         [(t_, 14, BONE, DISPLAY, False, 5, 0), (d, 10.5, STONE_D, TEXT, False, 0, 0)], line_spacing=1.3)
rect(s, M, Inches(5.35), Inches(11.53), Inches(0.72), WINE)
text(s, M + Inches(0.3), Inches(5.55), Inches(11), Inches(0.5),
     [("TARGETS — ", 9.5, PINK, TEXT, True, 0, 1.6),
      ("LCP under 2.5s on 4G · CLS under 0.1 · self-hosted fonts with metric-matched fallbacks · "
       "no autoplay with sound · semantic HTML, Article and Event schema.",
       11, BONE, TEXT, False, 0, 0)], line_spacing=1.4)
pagenum(s, 26, STONE_D)
notes(s, "Keep this short. The client is not buying a stack.",
      "Every choice traces back to something in the experience.",
      "'Here is how we would run it.'")
fade(s)

# ── 27 PROJECT PROCESS ────────────────────────────────────────────────────
s = blank(prs); bg(s, BONE)
eyebrow(s, M, Inches(0.75), "26 — Process")
text(s, M, Inches(1.15), Inches(9), Inches(0.7), [("Seven steps.", 34, INK, DISPLAY, False, 0, -0.9)])
steps2 = [("01", "Discover"), ("02", "Strategise"), ("03", "Wireframe"), ("04", "Design"),
          ("05", "Develop"), ("06", "Test"), ("07", "Launch")]
xx = M
for i, (n, t_) in enumerate(steps2):
    wd = Inches(1.5)
    c = WINE if i < 3 else (ROSE if i < 6 else INK)
    rect(s, xx, Inches(2.4), wd, Inches(1.5), c)
    text(s, xx, Inches(2.62), wd, Inches(0.4),
         [(n, 22, BONE, DISPLAY, False, 0, 0)], align=PP_ALIGN.CENTER)
    text(s, xx, Inches(3.25), wd, Inches(0.4),
         [(t_, 11.5, BONE, TEXT, True, 0, 0.8)], align=PP_ALIGN.CENTER)
    if i < 6:
        text(s, xx + wd, Inches(3.0), Inches(0.15), Inches(0.3),
             [("›", 15, STONE, TEXT, True, 0, 0)], align=PP_ALIGN.CENTER)
    xx += wd + Inches(0.15)
text(s, M, Inches(4.4), Inches(11.5), Inches(0.9),
     [("We do not skip the wireframe stage and we do not skip the concept comparison. "
       "Both already exist for SHUMI — this deck is the output of steps one to three.",
       13, STONE, TEXT, False, 0, 0)], line_spacing=1.5)
pagenum(s, 27)
notes(s, "Show process discipline.",
      "Steps 1–3 are already done and in front of them.",
      "'Which means the timeline starts from a running start.'")
fade(s)

# ── 28 TIMELINE ───────────────────────────────────────────────────────────
s = blank(prs); bg(s, BONE)
eyebrow(s, M, Inches(0.75), "27 — Timeline")
text(s, M, Inches(1.15), Inches(9), Inches(0.7), [("Seven to eleven weeks.", 34, INK, DISPLAY, False, 0, -0.9)])
phases = [("Discovery", 1, 1), ("UX / Wireframes", 1, 2), ("UI Design", 2, 2),
          ("Development", 2, 4), ("Content / Integrations", 1, 2), ("QA / Launch", 1, 1)]
maxw = 11.0; unit = maxw / 12.0
for i, (name, lo, hi) in enumerate(phases):
    yy = Inches(2.2) + Inches(0.62) * i
    text(s, M, yy, Inches(2.6), Inches(0.3), [(name, 12.5, INK, TEXT, False, 0, 0)])
    rect(s, M + Inches(2.7), yy + Inches(0.04), Inches(unit * lo), Inches(0.24), WINE)
    if hi > lo:
        rect(s, M + Inches(2.7) + Inches(unit * lo), yy + Inches(0.04),
             Inches(unit * (hi - lo)), Inches(0.24), ROSE)
    text(s, M + Inches(2.7) + Inches(unit * hi) + Inches(0.14), yy + Inches(0.03), Inches(1.6), Inches(0.3),
         [(f"{lo}" if lo == hi else f"{lo}–{hi} weeks", 10.5, STONE, TEXT, False, 0, 0)])
rect(s, M, Inches(6.05), Inches(11.53), Inches(0.7), INK)
text(s, M + Inches(0.3), Inches(6.25), Inches(11), Inches(0.4),
     [("THIS IS AN ESTIMATE. ", 9.5, PINK, TEXT, True, 0, 1.6),
      ("It moves with content availability, feedback cycles and final scope — content is almost always "
       "the critical path, not development.", 11, BONE, TEXT, False, 0, 0)])
pagenum(s, 28)
notes(s, "Be straight about the estimate.",
      "Content availability drives the timeline more than build does.",
      "'Here is exactly what you receive.'")
fade(s)

# ── 29 DELIVERABLES ───────────────────────────────────────────────────────
s = blank(prs); bg(s, INK)
eyebrow(s, M, Inches(0.8), "28 — Deliverables", PINK)
text(s, M, Inches(1.2), Inches(9), Inches(0.8), [("What you receive.", 34, BONE, DISPLAY, False, 0, -0.9)])
groups = [("Strategy", ["UX strategy", "Information architecture", "Three creative directions", "Wireframes"]),
          ("Design", ["Final UI design", "Responsive design", "Design system", "Motion specification"]),
          ("Build", ["Website development", "CMS", "Event architecture", "Forms"]),
          ("Launch", ["SEO foundation", "Analytics", "Accessibility audit", "QA, deployment & support"])]
for i, (g, items) in enumerate(groups):
    xx = M + Inches(2.9) * i
    rule(s, xx, Inches(2.3), Inches(2.6), PINK)
    text(s, xx, Inches(2.5), Inches(2.7), Inches(0.4), [(g, 19, BONE, DISPLAY, False, 0, 0)])
    for j, it in enumerate(items):
        yy = Inches(3.1) + Inches(0.46) * j
        rect(s, xx, yy + Inches(0.09), Inches(0.1), Pt(1.5), PINK)
        text(s, xx + Inches(0.26), yy, Inches(2.5), Inches(0.42),
             [(it, 11, BONE, TEXT, False, 0, 0)])
pagenum(s, 29, STONE_D)
notes(s, "Make the scope concrete.",
      "Sixteen deliverables across four stages.",
      "'And what it costs.'")
fade(s)

# ── 30 INVESTMENT ─────────────────────────────────────────────────────────
s = blank(prs); bg(s, BONE)
eyebrow(s, M, Inches(0.75), "29 — Investment")
text(s, M, Inches(1.15), Inches(9), Inches(0.7), [("Investment.", 34, INK, DISPLAY, False, 0, -0.9)])
inv = [("Strategy & UX", "Discovery, IA, wireframes, three directions"),
       ("UI / Visual design", "Final design, responsive, design system"),
       ("Development", "Frontend build, performance, accessibility"),
       ("CMS & integrations", "Content model, forms, event architecture"),
       ("Testing & launch", "QA, deployment, launch support")]
for i, (t_, d) in enumerate(inv):
    yy = Inches(2.2) + Inches(0.78) * i
    if i % 2 == 0: rect(s, M, yy - Inches(0.08), Inches(11.53), Inches(0.72), WHITE)
    text(s, M + Inches(0.25), yy, Inches(5.4), Inches(0.6),
         [(t_, 15, INK, DISPLAY, False, 3, 0), (d, 10, STONE, TEXT, False, 0, 0)], line_spacing=1.25)
    text(s, M + Inches(8.0), yy + Inches(0.05), Inches(3.2), Inches(0.4),
         [("$XX,XXX", 20, ROSE, DISPLAY, False, 0, -0.4)], align=PP_ALIGN.RIGHT)
yy = Inches(2.2) + Inches(0.78) * 5
rect(s, M, yy, Inches(11.53), Pt(2), INK)
text(s, M + Inches(0.25), yy + Inches(0.18), Inches(5.4), Inches(0.4),
     [("Total", 17, INK, DISPLAY, True, 0, 0)])
text(s, M + Inches(8.0), yy + Inches(0.16), Inches(3.2), Inches(0.4),
     [("$XX,XXX", 24, INK, DISPLAY, True, 0, -0.4)], align=PP_ALIGN.RIGHT)
rect(s, M, Inches(6.35), Inches(11.53), Inches(0.6), WINE)
text(s, M + Inches(0.3), Inches(6.52), Inches(11), Inches(0.3),
     [("Figures are placeholders. ", 10, PINK, TEXT, True, 0, 1.2),
      ("No pricing has been invented — these are filled in once scope is confirmed with SHUMI.",
       10.5, BONE, TEXT, False, 0, 0)])
pagenum(s, 30)
notes(s, "Present the shape of the investment without inventing a number.",
      "Five workstreams. Pricing follows scope confirmation.",
      "'A word on who is doing the work.'")
fade(s)

# ── 31 WHY TEDCANLABS ─────────────────────────────────────────────────────
s = blank(prs); bg(s, WINE)
p = img("c3-speaker.jpg")
if p: s.shapes.add_picture(p, Inches(8.0), 0, Inches(5.33), H)
scrim(s, Inches(8.0), 0, Inches(5.33), H, WINE, 0.32)
eyebrow(s, M, Inches(0.85), "30 — Why TedcanLabs", PINK)
text(s, M, Inches(1.45), Inches(6.6), Inches(0.9), [("Four things, done together.", 34, BONE, DISPLAY, False, 0, -0.9)])
why = [("Strategy", "We start with the argument, not the pages."),
       ("Design", "Premium digital experiences that hold up beside anyone's."),
       ("Technology", "Scalable implementation, built to be handed over."),
       ("Experience", "Human-centred products, measured rather than assumed.")]
for i, (t_, d) in enumerate(why):
    yy = Inches(2.6) + Inches(0.95) * i
    rect(s, M, yy + Inches(0.12), Inches(0.13), Pt(1.5), PINK)
    text(s, M + Inches(0.34), yy, Inches(6.2), Inches(0.8),
         [(t_, 19, BONE, DISPLAY, False, 4, 0), (d, 11.5, STONE_D, TEXT, False, 0, 0)], line_spacing=1.3)
text(s, M, Inches(6.45), Inches(6.6), Inches(0.5),
     [("The quality of this document is the argument.", 12, PINK, TEXT, False, 0, 0)])
pagenum(s, 31, STONE_D)
notes(s, "Keep it short and let the work speak.",
      "Do not oversell here — the previous thirty slides are the credential.",
      "'To close.'")
fade(s)

# ── 32 CLOSING ────────────────────────────────────────────────────────────
s = blank(prs); bg(s, INK)
has6 = video_full(s, "06_closing.mp4", dim=0.6)
if not has6 and not os.path.exists(os.path.join(VID, "06_closing_poster.jpg")):
    full_image(s, os.path.join(VID, "01_global_women_poster.jpg"))
    scrim(s, 0, 0, W, H, INK, 0.62)
text(s, M, Inches(1.75), Inches(10.2), Inches(2.4),
     [("SHUMI deserves a digital experience", 36, BONE, DISPLAY, False, 3, -1),
      ("that reflects the scale of its mission.", 36, BONE, DISPLAY, False, 0, -1)], line_spacing=1.05)
text(s, M, Inches(4.0), Inches(8.2), Inches(1.0),
     [("A platform that does not simply tell people what SHUMI does — but lets them experience "
       "the women, stories, communities and impact behind the organisation.",
       14, STONE_D, TEXT, False, 0, 0)], line_spacing=1.5)
rule(s, M, Inches(5.25), Inches(1.6), PINK, Pt(3))
text(s, M, Inches(5.55), Inches(9), Inches(0.9),
     [("Let's build SHUMI's digital future.", 30, PINK, DISPLAY, False, 0, -0.8)])
text(s, M, H - Inches(1.05), Inches(5), Inches(0.5),
     [("TEDCANLABS", 12, BONE, TEXT, True, 0, 2.4)])
pagenum(s, 32, STONE_D)
notes(s, "Close on scale and a clear next step.",
      "The ask is a decision on direction, not on budget.",
      "Stop talking. Let the room respond.")
fade(s, 1000)

# ── SAVE ──────────────────────────────────────────────────────────────────
out = os.path.join(ROOT, "PowerPoint", OUTNAME)
prs.save(out)
print(f"saved {out}")
print(f"slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")
